from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


class Create:
    def __init__(self):
        self.prs = Presentation()

    def create_prez(self, content):
        for slide in content:
            # Используем макет с заголовком и контентом
            slide_layout = self.prs.slide_layouts[1]  # Макет с заголовком и текстовым заполнителем
            slide_obj = self.prs.slides.add_slide(slide_layout)

            # Настройка фона слайда
            background = slide_obj.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(153, 153, 232)  # Темно-синий фон

            # Настройка заголовка
            title = slide_obj.shapes.title
            title.text = slide["title"]
            title.text_frame.paragraphs[0].font.size = Pt(28)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # Белый текст
            title.text_frame.paragraphs[0].font.name = "Arial"
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Настройка текстового контента
            content_placeholder = slide_obj.placeholders[1]
            content_placeholder.text = slide["content"]
            content_placeholder.width = Inches(4.5)  # Ограничиваем ширину текста
            content_placeholder.left = Inches(0.5)  # Сдвигаем текст влево
            content_placeholder.top = Inches(1.5)   # Отступ сверху
            for paragraph in content_placeholder.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)  # Белый текст
                paragraph.font.name = "Arial"
                paragraph.alignment = PP_ALIGN.LEFT

            # Добавление изображения, если путь существует
            img_path = slide.get("image", "1.jpg")  # По умолчанию '1.jpg'
            if os.path.exists(img_path):
                try:
                    # Размещаем изображение справа
                    left = Inches(5.5)  # Справа от текста
                    top = Inches(1.5)   # На одном уровне с текстом
                    width = Inches(4)   # Фиксированная ширина
                    slide_obj.shapes.add_picture(img_path, left, top, width=width)
                except Exception as e:
                    print(f"Ошибка при добавлении изображения {img_path}: {e}")
            else:
                print(f"Файл изображения {img_path} не найден")

        # Сохранение презентации
        try:
            self.prs.save('test.pptx')
            print("Презентация успешно сохранена как test.pptx")
        except Exception as e:
            print(f"Ошибка при сохранении презентации: {e}")

# Пример использования
if __name__ == "__main__":
    slides_content = [
        {
            "title": "Использование искусственного интеллекта (15.04.2025)",
            "content": "В данной презентации рассмотрим, как искусственный интеллект влияет на различные области нашей жизни.",
            "image": "1.jpg"
        },
        {
            "title": "Что такое Искусственный Интеллект?",
            "content": "Искусственный интеллект (ИИ) — это область компьютерной науки, которая сосредоточена на создании систем, способных выполнять задачи, требующие человеческого интеллекта.",
            "image": "1.jpg"
        }
    ]

    prs = Create()
    prs.create_prez(slides_content)